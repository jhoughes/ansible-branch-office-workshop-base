#!/usr/bin/env python3
"""
build-slides.py — Builds workshop.pptx using the conference template.

Run:
    python3 build-slides.py

Reads:  template.pptx (the conference's template with branding intact)
Writes: workshop.pptx

The conference template provides:
  - Slide size 13.33" x 7.5" (LAYOUT_WIDE)
  - Top-right "POWERSHELL DEVOPS 2026 GLOBAL SUMMIT" logo
  - Top-left "PowerShell + DevOps Global Summit" wordmark
  - CC-BY-SA license mark bottom-left
  - "April 13-16, 2026" date pill bottom-center
  - Slide number counter bottom-right
  - Vertical-line decorative backdrop
  - Title font: Space Grotesk, 44pt, dark navy

We use:
  - Layout 0 (Title Slide)            — slide 1, already populated by the template
  - Layout 14 (13_Custom Layout)      — section dividers (centered title only)
  - Layout 1 (Single content no pic)  — body slides (title + body placeholder)

Body placeholder on Layout 1: left=1.49" top=2.18" width=10.83" height=2.51"
Title placeholder on Layout 1: left=2.45" top=0.17" width=8.91" height=1.45"

Safe area: don't put content below ~6.7" (the bottom branding strip).
Don't put content above ~1.7" (the title placeholder).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# Conference palette
NAVY        = RGBColor(0x1A, 0x1A, 0x4A)
ACCENT      = RGBColor(0x29, 0xB6, 0xE1)   # the bright "Global Summit" blue
NEAR_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
MUTED       = RGBColor(0x60, 0x60, 0x70)
SOFT_BG     = RGBColor(0xF4, 0xF7, 0xFA)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG     = RGBColor(0x1A, 0x1A, 0x1A)
CODE_FG     = RGBColor(0xF0, 0xF0, 0xF0)
CODE_DIM    = RGBColor(0x88, 0x88, 0x88)
CODE_HI     = RGBColor(0x29, 0xB6, 0xE1)

# ---------------------------------------------------------------------------
# Setup
# ---------------------------------------------------------------------------
prs = Presentation("template.pptx")
LAYOUT_TITLE   = prs.slide_layouts[0]
LAYOUT_CONTENT = prs.slide_layouts[1]
LAYOUT_DIVIDER = prs.slide_layouts[14]


def fix_background(slide):
    """Inject the template's standard light-grey background override onto a
    newly-added slide. Without this, new slides render with the dark
    gradient that's the master's default, instead of the template's
    intended light grey."""
    from lxml import etree
    nsmap = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }
    bg_xml = """<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><p:bgPr><a:solidFill><a:schemeClr val="tx1"><a:lumMod val="85000"/></a:schemeClr></a:solidFill><a:effectLst/></p:bgPr></p:bg>"""
    bg = etree.fromstring(bg_xml)
    cSld = slide.element.find(qn('p:cSld'))
    # Insert as first child of cSld
    cSld.insert(0, bg)


def remove_template_extras():
    """Drop ALL of the template's starter slides (title, sponsor, thank-you).
    The template's title slide has a sizing bug that overflows; we'll
    build our own clean title slide. Sponsor isn't ours to populate.
    Thank-you we replace with our own."""
    sld_id_list = prs.slides._sldIdLst
    slide_ids = list(sld_id_list)
    for i in (2, 1, 0):
        rId = slide_ids[i].attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
        prs.part.drop_rel(rId)
        sld_id_list.remove(slide_ids[i])


def set_title(slide, text, size=36):
    title = slide.shapes.title
    if title is None:
        return
    title.text_frame.text = text
    for para in title.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Space Grotesk"
            run.font.size = Pt(size)


def suppress_bullet(para):
    """Override master's default bullet on a paragraph (use buNone)."""
    pPr = para._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    bu_none = pPr.makeelement(qn("a:buNone"), {})
    pPr.append(bu_none)


def hide_body_placeholder(slide):
    """Empty + suppress bullet on the body placeholder so we can use free
    textboxes instead. The placeholder still exists but renders nothing."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 10 and ph.has_text_frame:
            ph.text_frame.text = ""
            for para in ph.text_frame.paragraphs:
                suppress_bullet(para)


def add_textbox(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
                color=NEAR_BLACK, align="left", mono=False, name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    para = tf.paragraphs[0]
    if align == "center":
        para.alignment = PP_ALIGN.CENTER
    elif align == "right":
        para.alignment = PP_ALIGN.RIGHT
    suppress_bullet(para)
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Consolas" if mono else name
    return tb


def add_rect(slide, x, y, w, h, *, fill=None, line_color=None, line_w=1):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color is not None:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    return sh


def add_code_block_rich(slide, x, y, w, h, lines):
    """Dark monospace code block.
    lines = list. Each entry can be:
       - str (single fg-colored line, possibly empty)
       - list of (text, color) tuples (multi-run line)
    """
    add_rect(slide, x, y, w, h, fill=CODE_BG, line_color=NAVY, line_w=1)
    tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1),
                                   Inches(w - 0.3), Inches(h - 0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        suppress_bullet(para)
        if isinstance(line, str):
            runs = [(line, CODE_FG)]
        else:
            runs = line
        for text, color in runs:
            run = para.add_run()
            run.text = text
            run.font.name = "Consolas"
            run.font.size = Pt(13)
            run.font.color.rgb = color


# ===========================================================================
# Build deck
# ===========================================================================
remove_template_extras()
# (Template's broken title slide is gone; we build our own from scratch.)

# ---------------------------------------------------------------------------
# SLIDE 1: Our own title slide — uses LAYOUT_DIVIDER which renders cleanly
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_DIVIDER)
fix_background(s)
# The layout's centered title placeholder is at top=2.47" with 44pt font.
# Two-line title at 44pt fits comfortably; we position our subtitle and
# byline well below to avoid any overlap.
set_title(s, "Building End-to-End Automation\nwith Ansible", 44)

add_textbox(s, 1.5, 5.5, 10.3, 0.5,
            "A Hands-On Workshop for Solving Real-World Problems",
            size=16, italic=True, color=NAVY, align="center")
add_textbox(s, 1.5, 6.05, 10.3, 0.35,
            "Joe Houghes  ·  Mike Nelson  ·  @jhoughes  ·  @mnelson",
            size=12, color=MUTED, align="center")

# ---------------------------------------------------------------------------
# SLIDE 2: Welcome
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "Welcome")
hide_body_placeholder(s)

add_textbox(s, 1.5, 2.2, 5.0, 0.5, "Joe Houghes", size=20, bold=True, color=NAVY)
add_textbox(s, 1.5, 2.7, 5.0, 0.4, "Enterprise IT — Azure Stack, Cisco UCS, SCVMM",
            size=13, color=MUTED)
add_textbox(s, 1.5, 3.1, 5.0, 0.4, "Linux automation guy who lives in Microsoft shops",
            size=13, italic=True)

add_textbox(s, 7.0, 2.2, 5.3, 0.5, "Mike Nelson", size=20, bold=True, color=NAVY)
add_textbox(s, 7.0, 2.7, 5.3, 0.4, "Microsoft MVP — Cloud + Datacenter, PowerShell",
            size=13, color=MUTED)
add_textbox(s, 7.0, 3.1, 5.3, 0.4, "Windows admin who learned to love Ansible",
            size=13, italic=True)

add_rect(s, 1.5, 4.2, 10.3, 0.04, fill=ACCENT, line_color=ACCENT)
add_textbox(s, 1.5, 4.5, 10.3, 0.4, "Today's promise:", size=14, italic=True, color=MUTED)
add_textbox(s, 1.5, 4.95, 10.3, 0.7,
            "In 4 hours, you'll automate a hybrid Linux + Windows environment end-to-end.",
            size=20, bold=True, color=NAVY)

# ---------------------------------------------------------------------------
# SLIDE 3: The Scenario
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "The Scenario: New Branch Office Standup")
hide_body_placeholder(s)

add_textbox(s, 1.5, 2.1, 10.5, 0.5,
            "You just got the call. Corporate is opening a new branch office.",
            size=16, italic=True)
add_textbox(s, 1.5, 2.8, 10.5, 0.4, "You need:", size=16, bold=True, color=NAVY)

items = [
    "A Linux web server",
    "A Windows management host (admin tools, file share, users)",
    "Same security baseline applied to both",
    "Repeatable for the next 50 branches",
]
for i, item in enumerate(items):
    y = 3.4 + i * 0.55
    add_textbox(s, 1.7, y, 0.3, 0.5, "▸", size=18, bold=True, color=ACCENT)
    add_textbox(s, 2.1, y + 0.05, 9.5, 0.5, item, size=16)

# ---------------------------------------------------------------------------
# SLIDE 4: What you'll build (architecture)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "What You'll Build Today")
hide_body_placeholder(s)

def arch_box(x, label, sublabel, role_lines):
    add_rect(s, x, 2.5, 3.7, 3.2, fill=SOFT_BG, line_color=ACCENT, line_w=1.5)
    add_textbox(s, x, 2.65, 3.7, 0.5, label, size=22, bold=True, color=NAVY, align="center")
    add_textbox(s, x, 3.15, 3.7, 0.4, sublabel, size=11, italic=True, color=MUTED, align="center")
    for i, line in enumerate(role_lines):
        add_textbox(s, x + 0.2, 3.75 + i * 0.4, 3.3, 0.4, line, size=12, align="center")

arch_box(0.7, "control", "Ubuntu — public IP",
         ["Your Ansible control node", "The only host you SSH into", "Jump-box for everything"])
arch_box(4.8, "web1 + web2", "Ubuntu — private",
         ["nginx web tier", "Jinja2-templated content", "Rolling deploys in 3.3"])
arch_box(8.9, "mgmt1", "Windows Server 2022",
         ["Management host", "Chocolatey-installed tools", "Users, group, SMB share"])

add_textbox(s, 1.5, 5.95, 10.3, 0.4,
            "3 hosts per attendee  ·  Provisioned in Azure  ·  All you need is an SSH client",
            size=12, italic=True, color=MUTED, align="center")

# ---------------------------------------------------------------------------
# SLIDE 5: SECTION DIVIDER — Section 1
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_DIVIDER)
fix_background(s)
set_title(s, "1 — Foundations", 44)
add_textbox(s, 2.4, 4.1, 8.9, 0.5, "0:00 – 1:00",
            size=18, italic=True, color=ACCENT, align="center")

# ---------------------------------------------------------------------------
# SLIDE 6: Why Ansible
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "Why Ansible — Three Reasons")
hide_body_placeholder(s)

def reason(y, n, title, body):
    add_textbox(s, 1.5, y, 0.8, 0.9, n, size=44, bold=True, color=ACCENT)
    add_textbox(s, 2.5, y + 0.05, 9.5, 0.5, title, size=20, bold=True, color=NAVY)
    add_textbox(s, 2.5, y + 0.55, 9.5, 0.6, body, size=14)

reason(2.3, "1", "Agentless",
       "Pure SSH (Linux) and WinRM (Windows). Nothing to install on managed hosts.")
reason(3.7, "2", "Declarative + Idempotent",
       "You describe the desired state. Ansible figures out what to change.")
reason(5.1, "3", "Linux + Windows in one tool",
       "First-class Windows support, not an afterthought. You'll prove it today.")

# ---------------------------------------------------------------------------
# SLIDE 7: Your Lab Today
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "Your Lab Today")
hide_body_placeholder(s)

add_textbox(s, 1.5, 2.2, 10.5, 0.5,
            "Pick up your attendee card at the door.",
            size=20, bold=True, color=NAVY)
add_textbox(s, 1.5, 2.75, 10.5, 0.4,
            "It has your control node IP, your SSH password, and the Vault password.",
            size=13, italic=True)

# Repo box (left)
add_rect(s, 1.5, 3.55, 7.0, 2.5, fill=SOFT_BG, line_color=ACCENT, line_w=1)
add_textbox(s, 1.7, 3.7, 6.6, 0.4, "Workshop repo:", size=12, italic=True, color=MUTED)
add_textbox(s, 1.7, 4.15, 6.6, 0.5, "github.com/jhoughes/",
            size=18, bold=True, color=NAVY, mono=True)
add_textbox(s, 1.7, 4.65, 6.6, 0.5, "ansible-branch-office-workshop-base",
            size=18, bold=True, color=NAVY, mono=True)
add_textbox(s, 1.7, 5.45, 6.6, 0.4,
            "Already cloned at ~/workshop on your control node.",
            size=11, italic=True, color=MUTED)

# QR placeholder (right)
add_rect(s, 9.3, 3.55, 2.4, 2.4, fill=WHITE, line_color=NAVY, line_w=1)
add_textbox(s, 9.3, 4.55, 2.4, 0.4, "[QR CODE]", size=14, italic=True, color=MUTED, align="center")
add_textbox(s, 9.3, 6.05, 2.4, 0.3, "(replace with QR before workshop)",
            size=8, italic=True, color=MUTED, align="center")

# ---------------------------------------------------------------------------
# SLIDE 8: Ansible Fundamentals
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "Ansible Fundamentals — The Speed Run")
hide_body_placeholder(s)

def term(x, y, name, defn):
    add_textbox(s, x, y, 5.4, 0.4, name, size=14, bold=True, color=ACCENT, mono=True)
    add_textbox(s, x, y + 0.4, 5.4, 0.7, defn, size=12)

term(1.5, 2.3, "Inventory", "List of hosts and how to reach them. Static YAML or dynamic plugin.")
term(7.0, 2.3, "Playbook", "An ordered list of tasks against an ordered list of hosts.")
term(1.5, 3.7, "Module", "A unit of work: install a package, write a file, start a service.")
term(7.0, 3.7, "Role", "A reusable bundle of tasks, variables, templates, and handlers.")
term(1.5, 5.1, "Idempotency", "Run twice, get the same result. The thing that makes safe automation possible.")
term(7.0, 5.1, "Facts", "Information Ansible gathers about each host on connect.")

# ---------------------------------------------------------------------------
# SLIDE 9: SECTION DIVIDER — Section 2
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_DIVIDER)
fix_background(s)
set_title(s, "2 — Core Roles & Multi-Platform", 44)
add_textbox(s, 2.4, 4.1, 8.9, 0.5, "1:00 – 2:00",
            size=18, italic=True, color=ACCENT, align="center")

# ---------------------------------------------------------------------------
# SLIDE 10: 2.1 First Real Playbook
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "2.1 — First Real Playbook")
hide_body_placeholder(s)

add_textbox(s, 1.5, 2.2, 10.5, 0.5,
            "Goal: install nginx + deploy a Jinja2-templated status page on web1.",
            size=16, color=NAVY)

add_code_block_rich(s, 1.5, 3.0, 10.3, 1.5, [
    [("$ ", CODE_HI), ("cd ~/workshop", CODE_FG)],
    [("$ ", CODE_HI), ("ansible-playbook playbooks/02-web-tier.yml", CODE_FG)],
])

add_textbox(s, 1.5, 4.7, 10.3, 0.4,
            "Open LAB-2.1.md for the step-by-step.",
            size=13, italic=True, color=MUTED)

# Helper for code block lines (handles list of (text, color) tuples per line)
# Adjust the helper to support multi-run per line:
# We'll rewrite add_code_block above? It currently supports str or (text, color)
# but a single run per line. For our use we need multi-run lines, so:

# The code blocks with multi-run lines need a different helper. Replace
# the previous add_code_block call:

# (Done by writing a richer code block helper inline below)


def add_code_block_rich(slide, x, y, w, h, lines):
    """lines = list. Each entry can be:
       - str (single fg-colored line, possibly empty)
       - list of (text, color) tuples (multi-run line)
    """
    add_rect(slide, x, y, w, h, fill=CODE_BG, line_color=NAVY, line_w=1)
    tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1),
                                   Inches(w - 0.3), Inches(h - 0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        suppress_bullet(para)
        # Normalize to list of (text, color) tuples
        if isinstance(line, str):
            runs = [(line, CODE_FG)]
        else:
            runs = line
        for text, color in runs:
            run = para.add_run()
            run.text = text
            run.font.name = "Consolas"
            run.font.size = Pt(13)
            run.font.color.rgb = color


# Now redo slide 10's code block using the rich helper, replacing the simple one
# (Nuclear-simple approach: pop the simple code block off and add a rich one.)
# Actually the simple helper happens to work for the slide 10 case since each
# line is already a single-run-style list. So leave it.

# ---------------------------------------------------------------------------
# SLIDE 11: 2.2 Refactor into Roles
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "2.2 — Refactor Into Roles")
hide_body_placeholder(s)

add_textbox(s, 1.5, 2.2, 10.5, 0.5,
            "Same behavior. Dramatically better organization.",
            size=16, color=NAVY)

add_textbox(s, 1.5, 3.0, 4.8, 0.4, "Before (section 2.1):",
            size=12, bold=True, color=MUTED)
add_code_block_rich(s, 1.5, 3.45, 4.8, 2.0, [
    "playbooks/02-web-tier.yml",
    "",
    "~80 lines",
    "everything inline",
])

add_textbox(s, 7.0, 3.0, 4.8, 0.4, "After (section 2.2):",
            size=12, bold=True, color=ACCENT)
add_code_block_rich(s, 7.0, 3.45, 4.8, 2.0, [
    "playbooks/03-...-with-roles.yml",
    "",
    "~10 lines",
    "roles: [webserver]",
])

add_textbox(s, 1.5, 5.7, 10.3, 0.4,
            "Open LAB-2.2.md for the side-by-side.",
            size=13, italic=True, color=MUTED)

# ---------------------------------------------------------------------------
# SLIDE 12: 2.3 PowerShell + Windows
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "2.3 — PowerShell + Windows Automation")
hide_body_placeholder(s)

add_textbox(s, 1.5, 2.2, 10.5, 0.5,
            "Watch a Linux control node manage Windows.",
            size=16, italic=True, color=NAVY)

modules = [
    ("win_chocolatey", "install VS Code, Notepad++, 7-Zip, PowerShell 7"),
    ("win_user / win_group", "create alice + bob, BranchOfficeIT group"),
    ("win_share", "publish the BranchOfficeShared SMB share"),
    ("win_powershell", "run YOUR PowerShell script, capture structured output"),
]
for i, (mod, desc) in enumerate(modules):
    y = 3.0 + i * 0.75
    add_textbox(s, 1.7, y, 0.3, 0.4, "▸", size=16, bold=True, color=ACCENT)
    add_textbox(s, 2.1, y, 4.0, 0.4, mod, size=14, bold=True, color=ACCENT, mono=True)
    add_textbox(s, 6.2, y + 0.05, 5.7, 0.4, desc, size=13)

# ---------------------------------------------------------------------------
# SLIDE 13: SECTION DIVIDER — Section 3
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_DIVIDER)
fix_background(s)
set_title(s, "3 — Security & Orchestration", 44)
add_textbox(s, 2.4, 4.1, 8.9, 0.5, "2:10 – 3:10",
            size=18, italic=True, color=ACCENT, align="center")

# ---------------------------------------------------------------------------
# SLIDE 14: 3.1 Hardening Checklist
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "3.1 — Hardening as Code")
hide_body_placeholder(s)

add_textbox(s, 1.5, 2.2, 10.5, 0.5,
            "One role. Two operating systems. One 9-item baseline.",
            size=16, color=NAVY, italic=True)

items_3_1 = [
    ("1", "Secure communications", "(SSH config / RDP NLA)"),
    ("2", "Disable root remote login", ""),
    ("3", "Remove unused software", "(telnet / SMBv1)"),
    ("4", "Principle of least privilege", "(logs / lockout policy)"),
    ("5", "Automate updates", "(unattended-upgrades / Windows Update)"),
    ("6", "Configure a firewall", "(ufw / Windows Firewall)"),
    ("7", "Logging", "(rsyslog / Windows audit policy)"),
    ("8", "Monitor login attempts", "(fail2ban / password complexity)"),
    ("9", "Mandatory access control", "(AppArmor / UAC)"),
]
for i, (n, title, parens) in enumerate(items_3_1):
    col = 0 if i < 5 else 1
    row = i if i < 5 else i - 5
    x = 1.5 if col == 0 else 7.0
    y = 3.0 + row * 0.55
    add_textbox(s, x, y, 0.4, 0.4, n, size=14, bold=True, color=ACCENT)
    add_textbox(s, x + 0.45, y + 0.02, 4.7, 0.4, title, size=13, bold=True, color=NAVY)
    if parens:
        add_textbox(s, x + 0.45, y + 0.32, 4.7, 0.3, parens, size=10, italic=True, color=MUTED)

# ---------------------------------------------------------------------------
# SLIDE 15: 3.2 Vault
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "3.2 — Stop Putting Passwords in YAML")
hide_body_placeholder(s)

add_textbox(s, 1.5, 2.2, 10.5, 0.5,
            "You've been editing a Windows password into a plaintext file all morning.",
            size=14, italic=True, color=MUTED)
add_textbox(s, 1.5, 2.75, 10.5, 0.5, "Time to fix that.",
            size=20, bold=True, color=NAVY)

add_code_block_rich(s, 1.5, 3.6, 10.3, 2.2, [
    [("# inventory/group_vars/windows/vars.yml  (committed, plaintext)", CODE_DIM)],
    [("ansible_password: ", CODE_FG), ("\"{{ vault_windows_admin_password }}\"", CODE_HI)],
    "",
    [("# inventory/group_vars/windows/vault.yml  (encrypted, git-ignored)", CODE_DIM)],
    [("vault_windows_admin_password: ", CODE_FG), ("\"<the actual password>\"", CODE_HI)],
])

add_textbox(s, 1.5, 6.0, 10.3, 0.4,
            "Vault password for today: on your card AND on the next slide.",
            size=12, italic=True, color=MUTED)

# ---------------------------------------------------------------------------
# SLIDE 16: Vault Password (THE slide attendees screenshot)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "Today's Workshop Vault Password")
hide_body_placeholder(s)

add_textbox(s, 1.5, 2.5, 10.3, 0.5,
            "Put this in .vault-pass on your control node:",
            size=14, color=MUTED, italic=True, align="center")

add_rect(s, 2.5, 3.2, 8.3, 1.6, fill=WHITE, line_color=ACCENT, line_w=3)
add_textbox(s, 2.5, 3.55, 8.3, 0.9,
            "POWERSHELL&DEVOPS_SUMMIT_2026!",
            size=30, bold=True, color=NAVY, mono=True, align="center")
add_textbox(s, 2.5, 4.4, 8.3, 0.4,
            "(same password for everyone — also on your card)",
            size=11, italic=True, color=MUTED, align="center")

add_textbox(s, 1.5, 5.5, 10.3, 0.7,
            "In production: per-attendee passwords or external secret store "
            "(HashiCorp Vault, AWS Secrets Manager, Azure Key Vault).",
            size=11, italic=True, color=MUTED, align="center")

# ---------------------------------------------------------------------------
# SLIDE 17: 3.3 Rolling Deployments
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "3.3 — Rolling Deployments")
hide_body_placeholder(s)

add_textbox(s, 1.5, 2.2, 10.5, 0.5,
            "Deploy to one host at a time. Health-check before moving on.",
            size=16, color=NAVY)

add_code_block_rich(s, 1.5, 3.0, 10.3, 2.0, [
    "- name: Rolling deploy",
    "  hosts: webservers",
    [("  serial: 1                    ", CODE_HI), ("# one host at a time", CODE_DIM)],
    [("  max_fail_percentage: 0       ", CODE_HI), ("# stop on ANY failure", CODE_DIM)],
])

add_textbox(s, 1.5, 5.3, 10.3, 0.5,
            "pre_tasks: drain from LB  →  roles: [webserver]  →  post_tasks: health check + return",
            size=12, italic=True, align="center")

# ---------------------------------------------------------------------------
# SLIDE 18: SECTION DIVIDER — Section 4
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_DIVIDER)
fix_background(s)
set_title(s, "4 — Production Patterns", 44)
add_textbox(s, 2.4, 4.1, 8.9, 0.5, "3:10 – 4:00",
            size=18, italic=True, color=ACCENT, align="center")

# ---------------------------------------------------------------------------
# SLIDE 19: 4.1 Debugging Flags
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "4.1 — Testing & Debugging")
hide_body_placeholder(s)

add_textbox(s, 1.5, 2.2, 10.5, 0.5,
            "Five flags that solve >90% of broken playbooks:",
            size=16, color=NAVY)

flags = [
    ("--check --diff", "What WOULD this do? Show me the diffs."),
    ("-vvv", "Verbose enough to see SSH/WinRM details."),
    ("--start-at-task NAME", "Skip to where you got stuck."),
    ("--step", "Walk one task at a time, interactively."),
    ("--syntax-check", "Validate YAML and module names without running."),
]
for i, (flag, desc) in enumerate(flags):
    col = 0 if i < 3 else 1
    row = i if i < 3 else i - 3
    x = 1.5 if col == 0 else 7.0
    y = 3.0 + row * 1.05
    add_textbox(s, x, y, 5.0, 0.4, flag, size=14, bold=True, color=ACCENT, mono=True)
    add_textbox(s, x, y + 0.45, 5.0, 0.5, desc, size=12)

# ---------------------------------------------------------------------------
# SLIDE 20: 4.3 Capstone
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "4.3 — Capstone: site.yml")
hide_body_placeholder(s)

add_textbox(s, 1.5, 2.2, 10.5, 0.5,
            "One playbook. One command. Entire branch office, end to end.",
            size=16, italic=True, color=NAVY)

add_code_block_rich(s, 1.5, 3.0, 10.3, 2.4, [
    [("$ ", CODE_HI), ("ansible-playbook playbooks/site.yml", CODE_FG)],
    "",
    [("Phase 1: ", CODE_HI), ("webserver role  →  webservers (web1, web2)", CODE_FG)],
    [("Phase 2: ", CODE_HI), ("windows-mgmt role  →  mgmt1", CODE_FG)],
    [("Phase 3: ", CODE_HI), ("hardening role  →  branch_office (all hosts)", CODE_FG)],
])

add_textbox(s, 1.5, 5.6, 10.3, 0.5,
            "This is what you'd commit to git. This is what runs in CI. This is the goal.",
            size=13, italic=True, align="center")

# ---------------------------------------------------------------------------
# SLIDE 21: Resources / Where to go from here
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "Where to Go From Here")
hide_body_placeholder(s)

def res(x, y, label, items):
    add_textbox(s, x, y, 5.4, 0.4, label, size=14, bold=True, color=ACCENT)
    for i, item in enumerate(items):
        add_textbox(s, x + 0.2, y + 0.5 + i * 0.4, 5.0, 0.4, "• " + item, size=11)

res(1.5, 2.2, "Books", [
    "Ansible: Up and Running (3rd ed.) — Hochstein",
    "Ansible for DevOps — Geerling",
    "Practical Ansible — Sebbar/Ehgartner",
])
res(7.0, 2.2, "Video Series", [
    "Ansible 101 — Jeff Geerling (free, 15 ep.)",
    "Getting Started with Ansible — Learn Linux TV",
    "Conference talks linked in repo's README",
])
res(1.5, 4.5, "Community Roles", [
    "geerlingguy.* on Galaxy",
    "dev-sec.ansible-collection-hardening",
    "community.windows / chocolatey.chocolatey",
])
res(7.0, 4.5, "Production Steps", [
    "Run AWX or Ansible Automation Platform",
    "Switch to dynamic inventory",
    "Add Molecule tests",
    "External secret stores",
])

# ---------------------------------------------------------------------------
# SLIDE 22: Wrap-up / Thanks
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
fix_background(s)
set_title(s, "Thank You", 44)
hide_body_placeholder(s)

add_textbox(s, 1.5, 2.6, 10.3, 0.7,
            "You just built a hybrid Linux + Windows automated environment, end to end, in 4 hours.",
            size=18, italic=True, color=NAVY, align="center")

add_textbox(s, 1.5, 3.7, 10.3, 0.5,
            "Take the repo. Run the lab again at home. Use the patterns at work.",
            size=14, align="center")

add_rect(s, 2.5, 4.7, 8.3, 1.0, fill=SOFT_BG, line_color=ACCENT, line_w=2)
add_textbox(s, 2.5, 4.95, 8.3, 0.5,
            "github.com/jhoughes/ansible-branch-office-workshop-base",
            size=14, color=NAVY, mono=True, align="center")

add_textbox(s, 1.5, 6.05, 10.3, 0.4,
            "Joe Houghes  ·  Mike Nelson",
            size=12, color=MUTED, align="center")

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
prs.save("workshop.pptx")
print("Wrote workshop.pptx — " + str(len(prs.slides)) + " slides")
